import os
import re
import io
import json
import time
import asyncio
import hashlib
import traceback
from collections import Counter
from urllib.parse import urlparse

import httpx
from fastapi import FastAPI, Query, UploadFile, File, HTTPException
from fastapi.responses import HTMLResponse, JSONResponse, StreamingResponse
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from typing import Optional
from groq import Groq
from dotenv import load_dotenv

# Document parsing libs
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation

load_dotenv()

ANAKIN_API_KEY = os.getenv("ANAKIN_API_KEY")
GROQ_API_KEY = os.getenv("GROQ_API_KEY")

if not ANAKIN_API_KEY or not GROQ_API_KEY:
    raise RuntimeError(
        "Missing ANAKIN_API_KEY or GROQ_API_KEY. Make sure they're set in your .env file."
    )

groq_client = Groq(api_key=GROQ_API_KEY, timeout=60.0)
app = FastAPI(title="WireSavant", version="2.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# =============================================================================
# IN-MEMORY CACHE + COUNTERS
# =============================================================================
_cache: dict[str, dict] = {}          # url -> full result dict
_file_cache: dict[str, dict] = {}     # file hash -> result (avoid re-parsing duplicates)
_total_analyzed: int = 0               # counter of unique URLs analyzed
_total_files_analyzed: int = 0         # counter of unique files analyzed

MAX_FILE_BYTES = 15 * 1024 * 1024      # 15MB upload cap
ALLOWED_FILE_EXT = {".pdf", ".docx", ".pptx"}


# =============================================================================
# SCRAPING — Anakin async job pattern
# =============================================================================
async def scrape_with_anakin(url: str, use_browser: bool = True) -> str | None:
    headers = {
        "X-API-Key": ANAKIN_API_KEY,
        "Content-Type": "application/json",
    }
    payload = {
        "url": url,
        "country": "us",
        "useBrowser": use_browser,
        "generateJson": False,
    }

    try:
        async with httpx.AsyncClient(timeout=90.0, follow_redirects=True) as client:
            submit_resp = await client.post(
                "https://api.anakin.io/v1/url-scraper",
                headers=headers,
                json=payload,
            )
            print(f"Anakin submit: status={submit_resp.status_code}, body={submit_resp.text[:300]}")

            if submit_resp.status_code not in (200, 202):
                return None

            submit_data = submit_resp.json()
            job_id = submit_data.get("jobId") or submit_data.get("id")
            if not job_id:
                return None

            print(f"Anakin job submitted: {job_id}")

            for attempt in range(30):
                await asyncio.sleep(2)
                poll_resp = await client.get(
                    f"https://api.anakin.io/v1/url-scraper/{job_id}",
                    headers=headers,
                )
                poll_data = poll_resp.json()
                status = poll_data.get("status")
                print(f"Anakin poll #{attempt + 1}: status={status}")

                if status == "completed":
                    content = (
                        poll_data.get("markdown")
                        or poll_data.get("cleanedHtml")
                        or poll_data.get("html")
                    )
                    if content:
                        print(f"Anakin scraped {len(content)} chars")
                        return content[:8000]
                    return None

                if status == "failed":
                    return None

            return None

    except Exception as e:
        print(f"Anakin exception: {e}")
        traceback.print_exc()
        return None


async def scrape_fallback(url: str) -> str | None:
    """Raw HTTP GET fallback when Anakin fails."""
    try:
        headers = {
            "User-Agent": (
                "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
                "AppleWebKit/537.36 (KHTML, like Gecko) "
                "Chrome/120.0.0.0 Safari/537.36"
            )
        }
        async with httpx.AsyncClient(timeout=20.0, follow_redirects=True) as client:
            response = await client.get(url, headers=headers)

        html = response.text
        html = re.sub(r"<script[^>]*>.*?</script>", " ", html, flags=re.DOTALL | re.IGNORECASE)
        html = re.sub(r"<style[^>]*>.*?</style>", " ", html, flags=re.DOTALL | re.IGNORECASE)
        html = re.sub(r"<[^>]+>", " ", html)
        html = re.sub(r"\s{2,}", " ", html).strip()

        print(f"Fallback scrape: {len(html)} chars from {url}")
        return html[:5000] if html else None

    except Exception as e:
        print(f"Fallback scrape failed: {e}")
        return None


# =============================================================================
# OUTBOUND LINKS — extract from raw HTML
# =============================================================================
def extract_outbound_links(html_or_markdown: str, source_url: str, limit: int = 8) -> list[str]:
    """Pull external domains the page links out to. Cheap regex, no parser dep."""
    if not html_or_markdown:
        return []

    source_host = urlparse(source_url).netloc.lower().replace("www.", "")
    hrefs = re.findall(r'href=["\'](https?://[^"\']+)["\']', html_or_markdown, re.IGNORECASE)
    hrefs += re.findall(r'\]\((https?://[^)]+)\)', html_or_markdown, re.IGNORECASE)

    counter: Counter[str] = Counter()
    for href in hrefs:
        try:
            host = urlparse(href).netloc.lower().replace("www.", "")
            if not host or host == source_host:
                continue
            counter[host] += 1
        except Exception:
            continue

    return [host for host, _ in counter.most_common(limit)]


# =============================================================================
# WIRE — Anakin's structured-action layer (api.anakin.io/v1/wire)
# Opportunistically pulls structured data for sites Wire has a catalog for,
# instead of relying solely on raw scraping. Core workflow step, not optional —
# every URL analysis checks for a Wire match before/alongside scraping.
# =============================================================================
WIRE_BASE = "https://api.anakin.io/v1/wire"

# domain substring -> wire service slug
WIRE_SERVICES = {
    "npmjs.com": "npm",
    "pypi.org": "pypi",
    # Both verified against the live Anakin Wire catalog (active, real actions,
    # param names match what _guess_params provides). The others below were
    # guesses from marketing docs and have NOT been confirmed:
    #   "github.com": "github",        # confirmed DEAD — catalog has 0 actions
    #   "reddit.com": "reddit",         # unverified
    #   "stackoverflow.com": "stackoverflow",  # unverified
    #   "amazon.": "amazon",            # unverified
}
# NOTE: _guess_params() below already knows how to derive params for "reddit"
# and "stackoverflow" — wiring up either is a one-line addition to
# WIRE_SERVICES above *once verified live*. To verify against your own Anakin
# key, run:
#   curl -s https://api.anakin.io/v1/wire/catalog/reddit -H "X-API-Key: $ANAKIN_API_KEY"
#   curl -s https://api.anakin.io/v1/wire/catalog/stackoverflow -H "X-API-Key: $ANAKIN_API_KEY"
# If "actions" comes back non-empty, add the matching line (e.g.
# "reddit.com": "reddit") and it's live — no other code changes needed.
# I don't have outbound network access in this sandbox to run that check
# myself, so I've left these unverified rather than guess; whichever you
# confirm is a one-line flip.

_wire_catalog_cache: dict[str, list] = {}


def _match_wire_service(url: str) -> str | None:
    host = urlparse(url).netloc.lower()
    for needle, slug in WIRE_SERVICES.items():
        if needle in host:
            return slug
    return None


def _guess_params(slug: str, url: str) -> dict | None:
    """Best-effort param extraction from the URL path for common Wire services."""
    path = urlparse(url).path.strip("/")
    parts = [p for p in path.split("/") if p]

    if slug == "github" and len(parts) >= 2:
        return {"owner": parts[0], "repo": parts[1]}
    if slug == "npm" and parts:
        name = parts[-1] if parts[0] != "package" else parts[-1]
        return {"package": name, "package_name": name, "query": name}
    if slug == "pypi" and parts:
        name = parts[1] if parts[0] == "project" and len(parts) > 1 else parts[-1]
        return {"package": name, "package_name": name, "query": name}
    if slug == "reddit" and "r" in parts:
        idx = parts.index("r")
        if idx + 1 < len(parts):
            return {"subreddit": parts[idx + 1]}
    if slug == "stackoverflow" and "questions" in parts:
        idx = parts.index("questions")
        if idx + 1 < len(parts):
            return {"question_id": parts[idx + 1]}
    return None


async def wire_lookup(url: str) -> dict | None:
    """
    Core Wire workflow step. Returns a small dict of structured data for the
    URL if it matches a known Wire service, else None. Never raises — any
    failure (auth, catalog shape, missing params) just means no enrichment.
    """
    slug = _match_wire_service(url)
    if not slug:
        print(f"Wire: no service match for {url}")
        return None

    params = _guess_params(slug, url)
    if not params:
        print(f"Wire: matched service '{slug}' but couldn't derive params from {url}")
        return None

    return await _wire_call(slug, params)


async def wire_lookup_package(slug: str, package_name: str) -> dict | None:
    """
    Direct Wire lookup for the Compare feature — given an ecosystem slug
    ('npm' or 'pypi') and a bare package name, skips URL parsing entirely
    and calls Wire with the same param shape _guess_params would produce.
    """
    if slug not in ("npm", "pypi"):
        return None
    params = {"package": package_name, "package_name": package_name, "query": package_name}
    return await _wire_call(slug, params)


async def _wire_call(slug: str, params: dict) -> dict | None:
    """Shared Wire submit/poll logic used by both wire_lookup() and wire_lookup_package()."""
    print(f"Wire: matched service '{slug}', params={params}")
    headers = {"X-API-Key": ANAKIN_API_KEY, "Content-Type": "application/json"}

    try:
        async with httpx.AsyncClient(timeout=20.0) as client:
            # Discover actions for this service (cached per process)
            if slug in _wire_catalog_cache:
                actions = _wire_catalog_cache[slug]
            else:
                cat_resp = await client.get(f"{WIRE_BASE}/catalog/{slug}", headers=headers)
                print(f"Wire catalog '{slug}': status={cat_resp.status_code}, body={cat_resp.text[:500]}")
                if cat_resp.status_code != 200:
                    print(f"Wire catalog failed for '{slug}': {cat_resp.status_code}")
                    return None
                cat_data = cat_resp.json()
                actions = cat_data.get("actions")
                if not isinstance(actions, list):
                    print(f"Wire: catalog response for '{slug}' had unexpected shape: {cat_data!r}")
                    return None
                if not actions:
                    status = (cat_data.get("catalog") or {}).get("status")
                    print(f"Wire: '{slug}' catalog has 0 actions (catalog status='{status}') — nothing to call.")
                    return None
                _wire_catalog_cache[slug] = actions

            print(f"Wire: {len(actions)} action(s) available for '{slug}'")

            # Pick the action whose declared params best match what we extracted.
            # Each action's params live under "parameters": [{"name": ..., "required": ...}, ...]
            chosen_action_id = None
            chosen_call_params = None
            for action in actions:
                action_id = action.get("action_id") or action.get("id")
                declared = action.get("parameters") or []
                declared_names = {p.get("name") for p in declared if isinstance(p, dict) and p.get("name")}
                required_names = {p.get("name") for p in declared if isinstance(p, dict) and p.get("required") and p.get("name")}
                if required_names and required_names.issubset(params.keys()):
                    chosen_action_id = action_id
                    chosen_call_params = {k: v for k, v in params.items() if k in declared_names}
                    break
            if not chosen_action_id and actions:
                # fall back to the first listed action, passing through whatever we have
                first = actions[0]
                chosen_action_id = first.get("action_id") or first.get("id")
                declared = first.get("parameters") or []
                declared_names = {p.get("name") for p in declared if isinstance(p, dict) and p.get("name")}
                chosen_call_params = {k: v for k, v in params.items() if k in declared_names} or params
            if not chosen_action_id:
                print(f"Wire: no usable action_id found for '{slug}' (actions={actions!r})")
                return None

            print(f"Wire: calling action '{chosen_action_id}' with params={chosen_call_params}")
            task_resp = await client.post(
                f"{WIRE_BASE}/task",
                headers=headers,
                json={"action_id": chosen_action_id, "params": chosen_call_params},
            )
            print(f"Wire task '{chosen_action_id}': status={task_resp.status_code}, body={task_resp.text[:500]}")
            if task_resp.status_code not in (200, 202):
                print(f"Wire task failed for '{chosen_action_id}': {task_resp.status_code}")
                return None

            task_data = task_resp.json()

            # Actions are "mode":"async" — submit/poll pattern. The submit response
            # gives us the exact poll_url to use (e.g. "/v1/wire/jobs/{job_id}") —
            # use that instead of guessing the path.
            poll_url = task_data.get("poll_url") or task_data.get("pollUrl")
            status = task_data.get("status")
            data = task_data.get("data")

            if data is None and poll_url and status not in ("completed", "success"):
                if poll_url.startswith("http"):
                    full_poll_url = poll_url
                else:
                    full_poll_url = "https://api.anakin.io" + poll_url
                for attempt in range(25):
                    await asyncio.sleep(2)
                    poll_resp = await client.get(full_poll_url, headers=headers)
                    if poll_resp.status_code != 200:
                        print(f"Wire poll #{attempt + 1} failed: {poll_resp.status_code}")
                        continue
                    poll_data = poll_resp.json()
                    status = poll_data.get("status")
                    print(f"Wire poll #{attempt + 1} for '{chosen_action_id}': status={status}")
                    if status in ("completed", "success"):
                        data = poll_data.get("data") or poll_data.get("result")
                        break
                    if status in ("failed", "error"):
                        print(f"Wire task '{chosen_action_id}' failed during polling.")
                        return None

            if data is None:
                print(f"Wire: '{chosen_action_id}' produced no data (status={status})")
                return None

            # Anakin wraps the real payload in an envelope: {status, data, files, error, meta}
            if isinstance(data, dict) and "data" in data and isinstance(data.get("data"), dict):
                data = data["data"]

            print(f"Wire: enriched '{slug}' via {chosen_action_id}")
            return {"service": slug, "action_id": chosen_action_id, "data": data}

    except Exception as e:
        print(f"Wire lookup failed for '{slug}': {e}")
        traceback.print_exc()
        return None


# =============================================================================
# DOCUMENT TEXT EXTRACTION (PDF / DOCX / PPTX)
# =============================================================================
def extract_text_from_pdf(data: bytes) -> tuple[str, dict]:
    reader = PdfReader(io.BytesIO(data))
    chunks = []
    for page in reader.pages:
        try:
            chunks.append(page.extract_text() or "")
        except Exception:
            chunks.append("")
    text = "\n".join(chunks).strip()
    meta = {"pages": len(reader.pages), "type": "pdf"}
    return text, meta


def extract_text_from_docx(data: bytes) -> tuple[str, dict]:
    doc = DocxDocument(io.BytesIO(data))
    chunks = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text and cell.text.strip():
                    chunks.append(cell.text)
    text = "\n".join(chunks).strip()
    meta = {"paragraphs": len(doc.paragraphs), "type": "docx"}
    return text, meta


def extract_text_from_pptx(data: bytes) -> tuple[str, dict]:
    pres = Presentation(io.BytesIO(data))
    chunks = []
    for idx, slide in enumerate(pres.slides, start=1):
        slide_lines = [f"[Slide {idx}]"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_lines.append(shape.text)
        chunks.append("\n".join(slide_lines))
    text = "\n\n".join(chunks).strip()
    meta = {"slides": len(pres.slides), "type": "pptx"}
    return text, meta


def _extract_document_text_sync(filename: str, data: bytes) -> tuple[str, dict]:
    ext = os.path.splitext(filename.lower())[1]
    if ext not in ALLOWED_FILE_EXT:
        raise HTTPException(status_code=400, detail=f"Unsupported file type: {ext}. Allowed: {sorted(ALLOWED_FILE_EXT)}")
    if len(data) > MAX_FILE_BYTES:
        raise HTTPException(status_code=413, detail=f"File too large. Max {MAX_FILE_BYTES // (1024*1024)}MB.")

    if ext == ".pdf":
        text, meta = extract_text_from_pdf(data)
    elif ext == ".docx":
        text, meta = extract_text_from_docx(data)
    elif ext == ".pptx":
        text, meta = extract_text_from_pptx(data)
    else:
        # Should never reach here because of the check above, but keep it safe.
        raise HTTPException(status_code=400, detail=f"Unsupported file type: {ext}")

    if not text or len(text.strip()) < 20:
        raise HTTPException(status_code=422, detail="Could not extract readable text from this file. It may be image-only or corrupted.")
    return text[:8000], meta


async def extract_document_text(filename: str, data: bytes) -> tuple[str, dict]:
    """Async wrapper for the sync document parser so the event loop stays free."""
    return await asyncio.to_thread(_extract_document_text_sync, filename, data)


# =============================================================================
# SECURITY CHECK — returns risk_level: low / medium / high
# =============================================================================
def _check_security_with_ai_sync(source_label: str, content_excerpt: str, is_document: bool = False) -> tuple[bool, str, str]:
    """source_label is either a URL or a filename. Sync — wrap with asyncio.to_thread."""
    content_kind = "document" if is_document else "webpage"
    prompt = f"""Analyze this {content_kind} for phishing, scam, or security risks.

Source: {source_label}
Content excerpt: {content_excerpt[:2000]}

Evaluate:
- Does the {content_kind} try to harvest credentials, payment info, or personal data?
- Are there urgency / scare tactics ("Your account will be locked", "Act now or else")?
- Suspicious links, mismatched URLs, or requests to enable macros / external content?
- Generic greetings like "Dear Customer" paired with credential or payment requests?
- Typosquatting, suspicious short URLs, or impersonation of a known brand?
- For documents: are there requests to enable editing, macros, or external content?

IMPORTANT: Login forms, password fields, and credit card inputs on legitimate
websites (banks, SaaS products, e-commerce) are completely normal and NOT phishing.
Internal HR / finance documents that ask for routine info are also normal.

Respond with JSON only:
{{
  "is_safe": true,
  "reason": "brief explanation",
  "risk_level": "low"
}}
risk_level must be one of: "low", "medium", "high"
- low: no suspicious signals
- medium: some minor concerns but likely legitimate
- high: strong phishing or scam indicators"""

    completion = groq_client.chat.completions.create(
        messages=[{"role": "user", "content": prompt}],
        model="llama-3.1-8b-instant",
        response_format={"type": "json_object"},
    )
    raw = completion.choices[0].message.content
    result = json.loads(raw)
    is_safe = bool(result.get("is_safe", True))
    reason = result.get("reason") or "No issues detected"
    risk_level = result.get("risk_level") or "low"
    if risk_level not in ("low", "medium", "high"):
        risk_level = "low" if is_safe else "high"
    print(f"Security result: safe={is_safe}, risk={risk_level}, reason={reason}")
    return is_safe, str(reason), risk_level


async def check_security_with_ai(source_label: str, content_excerpt: str, is_document: bool = False) -> tuple[bool, str, str]:
    """Async wrapper that runs the sync Groq call in a thread so the event loop stays free."""
    try:
        return await asyncio.to_thread(_check_security_with_ai_sync, source_label, content_excerpt, is_document)
    except Exception as e:
        print(f"Security check failed: {e}")
        traceback.print_exc()
        return True, "Security check unavailable", "low"


# =============================================================================
# SUMMARIZATION + METADATA — expanded version (6 bullets + elevator pitch)
# =============================================================================
def _summarize_with_groq_sync(text: str, source_label: str, is_document: bool = False) -> dict:
    content_kind = "document" if is_document else "website"
    prompt = (
        f"You are an executive analyst. Read this {content_kind} content and respond ONLY with JSON.\n"
        "Extract:\n"
        "1. elevator_pitch: a single concise sentence (under 30 words) describing what this is.\n"
        "2. summary: array of exactly 6 strings (bullet points), each under 35 words, covering:\n"
        "   - What it does / what the company or document is about\n"
        "   - The target audience or intended readers\n"
        "   - Pricing model, cost, or monetization (if not mentioned, say 'Not specified')\n"
        "   - Technology stack, tools, or platforms detected\n"
        "   - Standout features, differentiators, or unique selling points\n"
        "   - Overall value proposition or key takeaway\n"
        "3. page_title: the apparent title, company name, or document title (string)\n"
        "4. tech_stack: array of detected technologies (e.g. React, Stripe, Cloudflare). Empty array if unknown.\n"
        "5. social_links: array of social platform names found (e.g. Twitter, LinkedIn). Empty array if none.\n"
        "6. content_type: one of 'product', 'service', 'documentation', 'presentation', 'report', 'marketing', 'other'\n\n"
        "Respond ONLY with JSON, no preamble:\n"
        '{"elevator_pitch": "...", "summary": [...], "page_title": "...", "tech_stack": [...], "social_links": [...], "content_type": "..."}\n\n'
        f"Source: {source_label}\nContent: {text[:3500]}"
    )
    completion = groq_client.chat.completions.create(
        messages=[{"role": "user", "content": prompt}],
        model="llama-3.1-8b-instant",
        response_format={"type": "json_object"},
    )
    raw = completion.choices[0].message.content
    result = json.loads(raw)
    # Normalize summary — handle None / scalar / list / nested-list from the model.
    # Groq occasionally returns one bullet as an actual JSON array (e.g. a tech-stack
    # bullet like ["Python", "HTTP", "Library"]) instead of a string. str()'ing that
    # directly leaks Python/JSON list syntax into the UI ("['Python', 'HTTP', ...]").
    # _flatten_bullet turns any such value into a clean, human-readable string.
    def _flatten_bullet(s) -> str:
        if isinstance(s, (list, tuple)):
            return ", ".join(_flatten_bullet(x) for x in s if x is not None)
        if isinstance(s, dict):
            return ", ".join(f"{k}: {_flatten_bullet(v)}" for k, v in s.items())
        return str(s).lstrip("- ").strip()

    summary_val = result.get("summary")
    if isinstance(summary_val, list):
        result["summary"] = [_flatten_bullet(s) for s in summary_val if s is not None]
    elif summary_val is None or summary_val == "":
        result["summary"] = []
    else:
        result["summary"] = [_flatten_bullet(summary_val)]
    # Override None values (setdefault leaves them in place)
    result["elevator_pitch"] = result.get("elevator_pitch") or ""
    result["page_title"] = result.get("page_title") or ""
    result["tech_stack"] = result.get("tech_stack") or []
    result["social_links"] = result.get("social_links") or []
    result["content_type"] = result.get("content_type") or "other"
    return result


async def summarize_with_groq(text: str, source_label: str, is_document: bool = False) -> dict:
    """Async wrapper for the sync Groq summary call."""
    try:
        return await asyncio.to_thread(_summarize_with_groq_sync, text, source_label, is_document)
    except Exception as e:
        print(f"Groq summary crashed: {e}")
        traceback.print_exc()
        return {
            "elevator_pitch": "",
            "summary": [f"Summary generation failed. Excerpt: {text[:200]}..."],
            "page_title": "",
            "tech_stack": [],
            "social_links": [],
            "content_type": "other",
        }


# =============================================================================
# ROUTES
# =============================================================================
@app.get("/", response_class=HTMLResponse)
async def read_root():
    with open("templates/index.html", "r", encoding="utf-8") as f:
        return f.read()


@app.get("/stats")
async def get_stats():
    return JSONResponse(content={
        "total_analyzed": _total_analyzed,
        "total_files_analyzed": _total_files_analyzed,
    })


class AnalyzeBody(BaseModel):
    url: str


class BatchBody(BaseModel):
    urls: list[str]


class CompareBody(BaseModel):
    ecosystem: str  # "npm" or "pypi"
    package_a: str
    package_b: str


@app.get("/analyze")
async def analyze_url_get(url: str = Query(...)):
    return await _analyze(url)


@app.post("/analyze")
async def analyze_url_post(body: AnalyzeBody):
    return await _analyze(body.url)


@app.get("/result")
async def result_page(url: str = Query(...), response_class=HTMLResponse):
    """Shareable result link — serves the same SPA; JS picks up ?url= on load."""
    with open("templates/index.html", "r", encoding="utf-8") as f:
        return HTMLResponse(content=f.read())


@app.post("/batch")
async def batch_analyze(body: BatchBody):
    urls = body.urls[:5]  # cap at 5
    tasks = [_analyze(u) for u in urls]
    results = await asyncio.gather(*tasks, return_exceptions=True)
    out = []
    for url, res in zip(urls, results):
        if isinstance(res, Exception):
            out.append({"url": url, "error": True, "summary": str(res)})
        else:
            data = res.body if hasattr(res, "body") else b"{}"
            out.append({"url": url, **json.loads(data)})
    return JSONResponse(content={"results": out})


@app.post("/compare")
async def compare_packages(body: CompareBody):
    """
    Compare mode — pull Wire data for two npm or pypi packages side-by-side.
    Makes Wire a visible, interactive feature instead of a quiet background step.
    """
    ecosystem = (body.ecosystem or "").strip().lower()
    if ecosystem not in ("npm", "pypi"):
        raise HTTPException(status_code=400, detail="ecosystem must be 'npm' or 'pypi'")

    name_a = (body.package_a or "").strip()
    name_b = (body.package_b or "").strip()
    if not name_a or not name_b:
        raise HTTPException(status_code=400, detail="Both package names are required.")

    results = await asyncio.gather(
        wire_lookup_package(ecosystem, name_a),
        wire_lookup_package(ecosystem, name_b),
        return_exceptions=True,
    )
    data_a, data_b = results

    def _clean(name, res):
        if isinstance(res, Exception) or res is None:
            return {"name": name, "found": False, "data": None}
        return {"name": name, "found": True, "data": res.get("data")}

    return JSONResponse(content={
        "ecosystem": ecosystem,
        "a": _clean(name_a, data_a),
        "b": _clean(name_b, data_b),
    })



class ExportDocxBody(BaseModel):
    title: str = "WireSavant Report"
    data: dict


def _build_report_docx(title: str, data: dict) -> io.BytesIO:
    """Render an analysis result dict into a styled .docx using python-docx."""
    from docx.shared import Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = DocxDocument()

    heading = doc.add_heading(title, level=0)
    for run in heading.runs:
        run.font.color.rgb = RGBColor(0x12, 0x1A, 0x1D)

    meta_p = doc.add_paragraph()
    meta_p.add_run(f"Generated {time.strftime('%Y-%m-%d %H:%M:%S UTC', time.gmtime())}").italic = True

    is_safe = data.get("is_safe")
    status = "Error" if data.get("error") else ("Clean" if is_safe else "Threat Detected")
    risk = data.get("risk_level")

    status_p = doc.add_paragraph()
    status_run = status_p.add_run(f"Status: {status}" + (f"  ·  Risk: {risk}" if risk else ""))
    status_run.bold = True
    status_run.font.size = Pt(13)
    if status == "Clean":
        status_run.font.color.rgb = RGBColor(0x1f, 0x8a, 0x4c)
    elif status == "Threat Detected":
        status_run.font.color.rgb = RGBColor(0xc0, 0x3a, 0x3a)
    else:
        status_run.font.color.rgb = RGBColor(0xb0, 0x7a, 0x1e)

    if data.get("page_title"):
        doc.add_paragraph(f"Title: {data['page_title']}")
    if data.get("filename"):
        doc.add_paragraph(f"File: {data['filename']}")
    if data.get("content_type") and data["content_type"] != "other":
        doc.add_paragraph(f"Content type: {data['content_type']}")

    if data.get("elevator_pitch") and is_safe:
        doc.add_heading("Elevator Pitch", level=2)
        p = doc.add_paragraph(data["elevator_pitch"])
        p.runs[0].italic = True

    summary = data.get("summary")
    if summary:
        doc.add_heading("Summary" if is_safe else "Risk Notes", level=2)
        bullets = summary if isinstance(summary, list) else [summary]
        for b in bullets:
            doc.add_paragraph(str(b), style="List Bullet")

    if data.get("tech_stack"):
        doc.add_heading("Tech Stack", level=2)
        doc.add_paragraph(", ".join(data["tech_stack"]))

    if data.get("social_links"):
        doc.add_heading("Social", level=2)
        doc.add_paragraph(", ".join(data["social_links"]))

    if data.get("outbound_links"):
        doc.add_heading("Outbound Links", level=2)
        doc.add_paragraph(", ".join(data["outbound_links"]))

    wire_data = data.get("wire_data")
    if wire_data:
        doc.add_heading("Wire Enrichment", level=2)
        doc.add_paragraph(f"Service: {wire_data.get('service', 'unknown')}  ·  Action: {wire_data.get('action_id', '')}")
        payload = wire_data.get("data")
        if isinstance(payload, dict):
            for k, v in payload.items():
                doc.add_paragraph(f"{k}: {v}", style="List Bullet")
        elif payload is not None:
            doc.add_paragraph(str(payload))

    footer_p = doc.add_paragraph()
    footer_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    footer_run = footer_p.add_run("Generated by WireSavant")
    footer_run.font.size = Pt(9)
    footer_run.font.color.rgb = RGBColor(0x8a, 0x8a, 0x8a)

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf


@app.post("/export-docx")
async def export_docx(body: ExportDocxBody):
    """Generate a downloadable Word (.docx) report from an analysis result."""
    try:
        buf = _build_report_docx(body.title, body.data)
    except Exception as e:
        print(f"DOCX export failed: {e}")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail="Could not generate report.")

    safe_name = re.sub(r"[^a-zA-Z0-9_\-.]+", "_", body.title)[:80] or "wiresavant_report"
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers={"Content-Disposition": f'attachment; filename="{safe_name}.docx"'},
    )


@app.post("/scan-file")
async def scan_file(file: UploadFile = File(...)):
    """Scan an uploaded document (PDF / DOCX / PPTX)."""
    global _total_files_analyzed

    filename = file.filename or "uploaded"
    ext = os.path.splitext(filename.lower())[1]
    if ext not in ALLOWED_FILE_EXT:
        raise HTTPException(status_code=400, detail=f"Unsupported file type '{ext}'. Allowed: {sorted(ALLOWED_FILE_EXT)}")

    data = await file.read()
    if not data:
        raise HTTPException(status_code=400, detail="Empty file.")
    if len(data) > MAX_FILE_BYTES:
        raise HTTPException(status_code=413, detail=f"File too large. Max {MAX_FILE_BYTES // (1024*1024)}MB.")

    # simple content fingerprint to dedupe (sha256 is fast; no need for thread offload)
    file_hash = hashlib.sha256(data).hexdigest()
    cache_key = f"file::{file_hash}"
    if cache_key in _file_cache:
        cached = dict(_file_cache[cache_key])
        cached["cached"] = True
        cached["filename"] = filename
        return JSONResponse(content=cached)

    t_start = time.time()
    try:
        text, doc_meta = await extract_document_text(filename, data)
    except HTTPException:
        raise
    except Exception as e:
        print(f"Document parse failed: {e}")
        traceback.print_exc()
        raise HTTPException(status_code=422, detail=f"Could not parse {ext.upper()} file. It may be corrupted or password-protected.")

    t_parse = time.time()
    print(f"Parsed {filename} → {len(text)} chars in {t_parse - t_start:.1f}s ({doc_meta})")

    # Security check
    try:
        is_safe, security_reason, risk_level = await check_security_with_ai(filename, text, is_document=True)
    except Exception as e:
        print(f"Security check crashed: {e}")
        is_safe, security_reason, risk_level = True, "", "low"

    t_ai = time.time()

    if not is_safe:
        result = {
            "is_safe": False,
            "summary": [security_reason],
            "risk_level": risk_level,
            "filename": filename,
            "file_meta": doc_meta,
            "parse_time": round(t_parse - t_start, 1),
            "ai_time": round(t_ai - t_parse, 1),
            "cached": False,
        }
        _file_cache[cache_key] = result
        _total_files_analyzed += 1
        return JSONResponse(content=result)

    # Summarize
    try:
        meta = await summarize_with_groq(text, filename, is_document=True)
        summary = meta.get("summary", [])
        elevator_pitch = meta.get("elevator_pitch", "")
        page_title = meta.get("page_title", "")
        tech_stack = meta.get("tech_stack", [])
        social_links = meta.get("social_links", [])
        content_type = meta.get("content_type", "other")
    except Exception as e:
        print(f"Groq summary crashed: {e}")
        traceback.print_exc()
        summary = [f"Document appears safe. Summary generation failed. Excerpt: {text[:200]}..."]
        elevator_pitch = ""
        page_title = ""
        tech_stack = []
        social_links = []
        content_type = "other"

    t_end = time.time()

    result = {
        "is_safe": True,
        "elevator_pitch": elevator_pitch,
        "summary": summary,
        "page_title": page_title,
        "tech_stack": tech_stack,
        "social_links": social_links,
        "content_type": content_type,
        "risk_level": risk_level,
        "filename": filename,
        "file_meta": doc_meta,
        "parse_time": round(t_parse - t_start, 1),
        "ai_time": round(t_end - t_parse, 1),
        "cached": False,
    }
    _file_cache[cache_key] = result
    _total_files_analyzed += 1
    return JSONResponse(content=result)


# =============================================================================
# CORE URL PIPELINE
# =============================================================================
async def _analyze(url: str):
    global _total_analyzed
    print(f"\n{'='*60}\nAnalyzing: {url}\n{'='*60}")

    if not url or not re.match(r"^https?://", url, re.IGNORECASE):
        return JSONResponse(content={
            "is_safe": False,
            "summary": ["Please provide a valid URL starting with http:// or https://"],
            "error": True,
        })

    # Cache hit
    if url in _cache:
        print(f"Cache hit for {url}")
        cached = dict(_cache[url])
        cached["cached"] = True
        return JSONResponse(content=cached)

    t_start = time.time()

    # Step 1: Scrape
    scraped_text = None
    try:
        scraped_text = await scrape_with_anakin(url, use_browser=True)
    except Exception as e:
        print(f"Anakin crashed: {e}")
        traceback.print_exc()

    # If Anakin returns minimal content, retry with longer wait
    if scraped_text and len(scraped_text) < 500:
        print(f"Anakin returned minimal content ({len(scraped_text)} chars), retrying...")
        retry = await scrape_with_anakin(url, use_browser=True)
        if retry and len(retry) > len(scraped_text):
            scraped_text = retry

    if not scraped_text:
        print("Anakin failed, trying fallback scraper...")
        try:
            scraped_text = await scrape_fallback(url)
        except Exception as e:
            print(f"Fallback also crashed: {e}")

    t_scrape = time.time()

    if not scraped_text:
        return JSONResponse(content={
            "is_safe": False,
            "summary": ["Could not load this page. It may be blocked, down, or too slow to respond."],
            "error": True,
        })

    print(f"Scraped {len(scraped_text)} characters in {t_scrape - t_start:.1f}s")

    # Outbound links (cheap, do before security check)
    outbound_links = extract_outbound_links(scraped_text, url)

    # Wire — structured-data enrichment for known services (core workflow step)
    wire_data = None
    try:
        wire_data = await wire_lookup(url)
    except Exception as e:
        print(f"Wire step crashed: {e}")

    # Step 2: Security check
    try:
        is_safe, security_reason, risk_level = await check_security_with_ai(url, scraped_text, is_document=False)
    except Exception as e:
        print(f"Security check crashed: {e}")
        is_safe = True
        security_reason = ""
        risk_level = "low"

    t_ai = time.time()

    if not is_safe:
        result = {
            "is_safe": False,
            "summary": [security_reason],
            "risk_level": risk_level,
            "outbound_links": outbound_links,
            "wire_data": wire_data,
            "scrape_time": round(t_scrape - t_start, 1),
            "ai_time": round(t_ai - t_scrape, 1),
            "cached": False,
        }
        _cache[url] = result
        _total_analyzed += 1
        return JSONResponse(content=result)

    # Step 3: Summarize + metadata
    try:
        meta = await summarize_with_groq(scraped_text, url, is_document=False)
        summary = meta.get("summary", [])
        elevator_pitch = meta.get("elevator_pitch", "")
        page_title = meta.get("page_title", "")
        tech_stack = meta.get("tech_stack", [])
        social_links = meta.get("social_links", [])
        content_type = meta.get("content_type", "other")
        print(f"Summary generated in {time.time() - t_ai:.1f}s")
    except Exception as e:
        print(f"Groq summary crashed: {e}")
        traceback.print_exc()
        summary = [f"Page appears safe. Summary generation failed. Excerpt: {scraped_text[:200]}..."]
        elevator_pitch = ""
        page_title = ""
        tech_stack = []
        social_links = []
        content_type = "other"

    t_end = time.time()

    result = {
        "is_safe": True,
        "elevator_pitch": elevator_pitch,
        "summary": summary,
        "page_title": page_title,
        "tech_stack": tech_stack,
        "social_links": social_links,
        "content_type": content_type,
        "outbound_links": outbound_links,
        "wire_data": wire_data,
        "risk_level": risk_level,
        "scrape_time": round(t_scrape - t_start, 1),
        "ai_time": round(t_end - t_scrape, 1),
        "cached": False,
    }
    _cache[url] = result
    _total_analyzed += 1
    return JSONResponse(content=result)


if __name__ == "__main__":
    import uvicorn
    uvicorn.run("main:app", host="0.0.0.0", port=8000, reload=True)